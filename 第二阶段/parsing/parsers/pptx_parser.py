"""PPTX 基础文本解析器。"""

from pathlib import Path

from 第二阶段.parsing.base import BaseParser, ParserDependencyError, ParserError
from 第二阶段.schemas.models import ParsedDocument


class PptxParser(BaseParser):
    file_type = "pptx"

    def parse(self, file_path: str | Path) -> ParsedDocument:
        path = self._validate_path(file_path)
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserDependencyError("PPTX 解析需要安装 python-pptx") from exc
        try:
            presentation = Presentation(str(path))
            blocks: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                blocks.append(f"[Slide {index}]")
                blocks.extend(
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
        except Exception as exc:
            raise ParserError(f"PPTX 解析失败：{path}") from exc
        return self._build_document(
            path, "\n".join(blocks), metadata={"slide_count": len(presentation.slides)}
        )

